from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 创建输出目录
output_dir = 'c:/Users/frankechen/CodeBuddy/chrome/AI_规划_2026'
os.makedirs(output_dir, exist_ok=True)

# 颜色定义
PRIMARY = RGBColor(0x1a, 0x36, 0x5d)      # 深蓝
SECONDARY = RGBColor(0x2c, 0x52, 0x82)    # 中蓝
ACCENT = RGBColor(0x31, 0x82, 0xce)       # 亮蓝
SUCCESS = RGBColor(0x38, 0xa1, 0x69)      # 绿
WARNING = RGBColor(0xd6, 0x9e, 0x2e)      # 黄
DANGER = RGBColor(0xe5, 0x3e, 0x3e)       # 红
LIGHT = RGBColor(0xf7, 0xfa, 0xfc)        # 浅灰
DARK = RGBColor(0x1a, 0x20, 0x2c)         # 深灰
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xec, 0xc9, 0x4b)         # 金色

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 分隔线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(3.7), Inches(3), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # 底部信息
    info_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2026年度规划 | 政企客户成功中心"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, title, number):
    """添加章节分隔页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY
    background.line.fill.background()
    
    # 章节编号
    num_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = LIGHT
    background.line.fill.background()
    
    # 标题栏
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.fill.background()
    
    # 标题文字
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 内容区域
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)
    
    return slide

def create_presentation():
    """创建完整的AI战略规划PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9比例
    
    print("开始创建AI年度规划PPT...")
    
    # ===== 封面 =====
    print("创建封面...")
    add_title_slide(prs, "AI赋能·智领未来", "政企客户成功中心 2026年度AI战略规划")
    
    # ===== 目录 =====
    print("创建目录...")
    add_content_slide(prs, "目录", [
        "01 五看分析：洞察趋势与机会",
        "02 三定规划：明确战略与目标",
        "03 战略屋：整体架构设计",
        "04 小组规划框架与目标",
        "05 实施路径与里程碑",
        "06 资源保障与预期成效"
    ])
    
    # ===== 第一部分：五看分析 =====
    print("创建五看分析...")
    add_section_slide(prs, "五看分析", "01")
    
    add_content_slide(prs, "五看分析框架", [
        "📈 看趋势：AI技术发展趋势、行业变革方向",
        "🎯 看市场：客户需求变化、市场机会窗口",
        "⚔️ 看竞争：竞争对手布局、行业最佳实践",
        "🔍 看自己：内部能力评估、优势与短板",
        "💡 看机会：战略机会点、破局切入点"
    ])
    
    add_content_slide(prs, "看趋势：2026 AI技术发展趋势", [
        "🚀 技术演进趋势：",
        "   • Agentic AI时代：从对话助手到自主智能体",
        "   • 多模态融合：文本+图像+语音+代码统一处理",
        "   • 上下文工程：长上下文、RAG成为标配",
        "   • 模型即服务：推理成本持续下降",
        "",
        "🏢 产业应用趋势：",
        "   • 数字分身/员工：7×24小时AI工作伙伴",
        "   • 个人AI助理：Clawbot爆火，Agentic工作流",
        "   • 复利工程：知识沉淀成为核心竞争力"
    ])
    
    add_content_slide(prs, "看市场/客户：需求洞察与痛点分析", [
        "⚠️ 核心痛点：",
        "   • 合同评审效率瓶颈：人工2-3天/份，风险条款易遗漏",
        "   • 项目清淤难度大：历史项目积压，人工梳理耗时长",
        "   • 出海合规复杂度高：多国法规复杂，决策周期6+月",
        "   • 知识传承断层：经验分散难沉淀，培训成本高",
        "",
        "✅ AI赋能机会点：",
        "   • 智能合同评审与风险识别",
        "   • 项目健康度智能诊断",
        "   • 出海合规知识库与问答",
        "   • 客户成功智能助手"
    ])
    
    add_content_slide(prs, "看竞争：内外部AI能力建设对标", [
        "🏢 内部竞争格局：",
        "   • 腾讯研究院-晓辉博士：复利工程、FDE前沿部署工程师",
        "   • 余一-AIasme：数字分身技术探索",
        "   • 安灯、云一、TEG等团队：研发场景AI提效实践",
        "",
        "🌐 外部对标分析：",
        "   • 公司AI专家Jeff：AI领导力进化论，超级个体到超级军团",
        "   • 京东数字员工、数字军团",
        "   • Clawbot个人AI助理",
        "   • Anthropic Claude Code工程实践"
    ])
    
    add_content_slide(prs, "看自己：现状盘点与能力评估", [
        "🏆 已有成果（2025）：",
        "   • 13个AI智能体，8个已上线",
        "   • 50+项目评审，60+用户数",
        "   • 预计年节省2000万",
        "   • 4篇KM知识沉淀",
        "",
        "⚠️ 差距与提升空间：",
        "   • Agent自主化程度有待提升，当前多为辅助型",
        "   • 知识沉淀体系化不足，经验复用率需提高",
        "   • 跨团队协作深度有限，复利效应未充分释放",
        "   • 个人AI助理覆盖率低，规模化推广待加强"
    ])
    
    add_content_slide(prs, "看机会：战略机会点识别", [
        "🎯 机会点1：Agent规模化（高优先级）",
        "   从13个智能体扩展到全员覆盖，打造人均AI助理工作模式",
        "",
        "🔄 机会点2：复利工程体系（核心战略）",
        "   构建知识沉淀-复用-进化的飞轮，实现边际成本递减",
        "",
        "🤖 机会点3：数字分身探索（创新突破）",
        "   探索数字员工、个人AI助理落地，构建7×24小时服务能力",
        "",
        "🌐 机会点4：生态协同（战略合作）",
        "   与研究院、TEG等团队深度协同，形成组织级AI能力网络"
    ])
    
    # ===== 第二部分：三定规划 =====
    print("创建三定规划...")
    add_section_slide(prs, "三定规划", "02")
    
    add_content_slide(prs, "三定规划框架", [
        "定战略：",
        "   • Agentic AI战略定位",
        "   • 复利工程核心路径",
        "   • 组织AI能力建设",
        "",
        "定目标：",
        "   • 量化业务指标",
        "   • 能力建设目标",
        "   • 覆盖率与渗透度",
        "",
        "定策略：",
        "   • 实施路径规划",
        "   • 关键举措分解",
        "   • 资源保障机制"
    ])
    
    add_content_slide(prs, "定战略：AI赋能战略定位", [
        "🎯 战略定位：",
        "   • 愿景：成为公司AI赋能标杆团队",
        "   • 使命：让AI成为每个人的超级助手",
        "   • 价值观：复利思维、持续进化、开放协同",
        "",
        "🔄 复利工程战略：",
        "   • 知识沉淀：每次项目经验编码为可复用资产",
        "   • 经验复用：通过Context工程实现知识自动加载",
        "   • 持续进化：建立反馈闭环，越用越聪明",
        "",
        "三大核心方向：Agent规模化、复利工程体系、数字分身探索"
    ])
    
    add_content_slide(prs, "定目标：2026年度核心目标", [
        "📊 业务指标目标：",
        "   • AI智能体数量：13 → 50+",
        "   • 活跃用户覆盖率：60 → 100人",
        "   • 预计年节省成本：2000万 → 5000万",
        "   • 效率提升幅度：平均50%+",
        "",
        "🎯 能力建设目标：",
        "   • 知识库覆盖：70MB → 200MB",
        "   • KM沉淀文章：4 → 20篇",
        "   • 培训覆盖人次：200 → 500人",
        "   • Agent自主化率：20% → 60%"
    ])
    
    add_content_slide(prs, "定策略：四大实施策略", [
        "🎯 策略一：Agent规模化",
        "   梳理全业务流程，按优先级分批次开发部署Agent",
        "",
        "🔄 策略二：复利工程",
        "   建立Context规范，每次项目沉淀可复用经验",
        "",
        "🤖 策略三：数字分身",
        "   选择高频标准化场景试点，探索7×24小时服务",
        "",
        "🌐 策略四：生态协同",
        "   与研究院、TEG建立复利工程协作，输出最佳实践"
    ])
    
    # ===== 第三部分：战略屋与实施 =====
    print("创建战略屋与实施...")
    add_section_slide(prs, "战略屋", "03")
    
    add_content_slide(prs, "战略屋：整体架构设计", [
        "🏆 愿景：成为公司AI赋能标杆团队",
        "",
        "三大核心支柱：",
        "   🎯 Agent规模化：50+智能体覆盖、全场景AI赋能",
        "   🔄 复利工程：知识沉淀体系、经验复用机制",
        "   🤖 数字分身：数字员工试点、7×24服务",
        "",
        "🏗️ 支撑体系：技术底座、人才体系、知识管理、生态协同"
    ])
    
    add_content_slide(prs, "小组规划框架：各组目标与举措", [
        "🎯 业务组A：",
        "   目标：合同评审AI覆盖率100%，效率提升80%",
        "   举措：鹰眼智审优化、新场景Agent开发",
        "",
        "🎯 业务组B：",
        "   目标：项目AI问诊覆盖200+项目，损益量化500万",
        "   举措：问诊Agent能力提升、历史项目数据治理",
        "",
        "🎯 业务组C：",
        "   目标：出海合规服务100+用户，决策效率提升60%",
        "   举措：知识库扩容、多语言支持增强"
    ])
    
    add_content_slide(prs, "实施路径：2026年度路线图", [
        "Q1 规划启动：",
        "   年度规划评审、复利工程体系搭建、Agent需求梳理",
        "",
        "Q2 规模推广：",
        "   30个Agent上线、核心场景全覆盖、知识库扩容",
        "",
        "Q3 创新突破：",
        "   数字分身试点、Agent自主化率提升、生态协同深化",
        "",
        "Q4 总结升华：",
        "   年度成果盘点、方法论沉淀输出、2027规划启动"
    ])
    
    add_content_slide(prs, "资源保障：支撑体系", [
        "👥 人才保障：",
        "   AI赋能负责人、各小组AI联络员、全员AI素养培训",
        "",
        "💰 预算保障：",
        "   大模型API费用、Dify平台订阅、知识库建设",
        "",
        "🛠️ 技术保障：",
        "   Dify平台稳定运行、RAG知识库、MCP工具集成",
        "",
        "📋 机制保障：",
        "   月度复盘、季度评审、知识沉淀、激励表彰"
    ])
    
    add_content_slide(prs, "预期成效：价值与影响", [
        "📊 业务成效：",
        "   效率提升平均50%+、成本节省5000万/年",
        "   风险识别准确率92%+、响应速度分钟级",
        "",
        "👥 组织成效：",
        "   全员AI素养提升、知识沉淀体系建立",
        "   团队协作效率优化、创新氛围增强",
        "",
        "🏆 品牌成效：",
        "   成为公司AI赋能标杆团队、方法论可复用推广"
    ])
    
    # ===== 结束页 =====
    print("创建结束页...")
    add_title_slide(prs, "AI赋能·智领未来", "让AI成为每个人的超级助手")
    
    # 保存文件
    output_path = os.path.join(output_dir, 'AI战略规划_2026年度完整版.pptx')
    prs.save(output_path)
    print(f"\nPPT创建完成！")
    print(f"保存位置：{output_path}")
    print(f"共 {len(prs.slides)} 页幻灯片")

if __name__ == '__main__':
    create_presentation()
